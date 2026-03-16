#!/usr/bin/env python3
"""
extract_doc.py - 从 PDF/PPT/Word 文档提取结构化文本内容

用法：
  python extract_doc.py <文件路径> [--output <输出JSON路径>] [--max-chars 50000]

输出 JSON 格式：
{
  "file_type": "pdf|pptx|docx",
  "title": "文档标题",
  "total_pages": 10,
  "sections": [
    {"index": 1, "title": "章节标题或页码", "content": "正文内容"}
  ],
  "full_text": "合并后的全文（用于分析）",
  "metadata": {"author": "", "created": ""}
}
"""

import sys
import json
import argparse
from pathlib import Path


def extract_pdf(filepath: str, max_chars: int) -> dict:
    try:
        import fitz  # pymupdf
    except ImportError:
        print("需要安装 pymupdf：pip install pymupdf", file=sys.stderr)
        sys.exit(1)

    doc = fitz.open(filepath)
    sections = []
    full_parts = []

    for i, page in enumerate(doc):
        text = page.get_text("text").strip()
        if not text:
            continue
        # 用每页的第一行作为标题（通常是章节名或页眉）
        lines = text.splitlines()
        title = lines[0].strip() if lines else f"第 {i+1} 页"
        sections.append({
            "index": i + 1,
            "title": title,
            "content": text
        })
        full_parts.append(text)

    full_text = "\n\n".join(full_parts)
    if len(full_text) > max_chars:
        full_text = full_text[:max_chars] + f"\n\n[内容已截断，原始长度 {len(''.join(full_parts))} 字符]"

    meta = doc.metadata or {}
    return {
        "file_type": "pdf",
        "title": meta.get("title") or Path(filepath).stem,
        "total_pages": len(doc),
        "sections": sections,
        "full_text": full_text,
        "metadata": {
            "author": meta.get("author", ""),
            "created": meta.get("creationDate", "")
        }
    }


def extract_pptx(filepath: str, max_chars: int) -> dict:
    try:
        from pptx import Presentation
    except ImportError:
        print("需要安装 python-pptx：pip install python-pptx", file=sys.stderr)
        sys.exit(1)

    prs = Presentation(filepath)
    sections = []
    full_parts = []

    for i, slide in enumerate(prs.slides):
        texts = []
        title = f"第 {i+1} 页"
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            shape_text = shape.text_frame.text.strip()
            if not shape_text:
                continue
            # 优先取标题占位符作为标题
            if shape.shape_type == 13 or (hasattr(shape, 'placeholder_format') and
                    shape.placeholder_format and shape.placeholder_format.idx == 0):
                title = shape_text.splitlines()[0]
            texts.append(shape_text)

        content = "\n".join(texts)
        if content:
            sections.append({"index": i + 1, "title": title, "content": content})
            full_parts.append(content)

    full_text = "\n\n".join(full_parts)
    if len(full_text) > max_chars:
        full_text = full_text[:max_chars] + f"\n\n[内容已截断]"

    core = prs.core_properties
    return {
        "file_type": "pptx",
        "title": (core.title or Path(filepath).stem),
        "total_pages": len(prs.slides),
        "sections": sections,
        "full_text": full_text,
        "metadata": {
            "author": core.author or "",
            "created": str(core.created) if core.created else ""
        }
    }


def extract_docx(filepath: str, max_chars: int) -> dict:
    try:
        import docx
    except ImportError:
        print("需要安装 python-docx：pip install python-docx", file=sys.stderr)
        sys.exit(1)

    doc = docx.Document(filepath)
    sections = []
    full_parts = []
    current_section = None
    current_content = []
    section_idx = 0

    for para in doc.paragraphs:
        text = para.text.strip()
        if not text:
            continue

        is_heading = para.style.name.startswith("Heading")
        if is_heading:
            # 保存上一个段落
            if current_section is not None:
                content = "\n".join(current_content)
                sections.append({
                    "index": section_idx,
                    "title": current_section,
                    "content": content
                })
                full_parts.append(content)
            section_idx += 1
            current_section = text
            current_content = []
        else:
            current_content.append(text)

    # 保存最后一个段落
    if current_section is not None:
        content = "\n".join(current_content)
        sections.append({
            "index": section_idx,
            "title": current_section,
            "content": content
        })
        full_parts.append(content)
    elif current_content:
        # 没有标题的文档，整体作为一个 section
        content = "\n".join(current_content)
        sections.append({"index": 1, "title": "正文", "content": content})
        full_parts.append(content)

    full_text = "\n\n".join(full_parts)
    if len(full_text) > max_chars:
        full_text = full_text[:max_chars] + f"\n\n[内容已截断]"

    core = doc.core_properties
    return {
        "file_type": "docx",
        "title": core.title or Path(filepath).stem,
        "total_pages": len(sections),
        "sections": sections,
        "full_text": full_text,
        "metadata": {
            "author": core.author or "",
            "created": str(core.created) if core.created else ""
        }
    }


def main():
    parser = argparse.ArgumentParser(description="从文档提取结构化文本")
    parser.add_argument("filepath", help="文件路径（PDF/PPTX/DOCX）")
    parser.add_argument("--output", "-o", help="输出 JSON 文件路径（默认输出到 stdout）")
    parser.add_argument("--max-chars", type=int, default=60000,
                        help="全文最大字符数（默认 60000）")
    args = parser.parse_args()

    filepath = args.filepath
    if not Path(filepath).exists():
        print(f"错误：文件不存在：{filepath}", file=sys.stderr)
        sys.exit(1)

    suffix = Path(filepath).suffix.lower()
    extractors = {
        ".pdf": extract_pdf,
        ".pptx": extract_pptx,
        ".ppt": extract_pptx,
        ".docx": extract_docx,
        ".doc": extract_docx,
    }

    if suffix not in extractors:
        print(f"错误：不支持的文件类型 {suffix}，支持 pdf/pptx/docx", file=sys.stderr)
        sys.exit(1)

    result = extractors[suffix](filepath, args.max_chars)
    result["source_file"] = filepath

    output_json = json.dumps(result, ensure_ascii=False, indent=2)
    if args.output:
        Path(args.output).parent.mkdir(parents=True, exist_ok=True)
        Path(args.output).write_text(output_json, encoding="utf-8")
        print(f"已提取到：{args.output}（{result['total_pages']} 页/章节，"
              f"{len(result['full_text'])} 字符）")
    else:
        print(output_json)


if __name__ == "__main__":
    main()
